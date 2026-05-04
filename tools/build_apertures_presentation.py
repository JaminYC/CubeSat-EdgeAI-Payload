"""
Genera presentacion del experimento de mascaras opticas usando el template INTISAT.

Output: Documentos de Referencia/Presentacion_Mascaras_Lensless.pptx

Ejecutar:
    python tools/build_apertures_presentation.py
"""

import os
import sys
import io
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# UTF-8 stdout para Windows
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "Documentos de Referencia" / "INTISAT_presentation_template[NO_MODIFICAR].pptx"
OUT_PATH = ROOT / "Documentos de Referencia" / "Presentacion_Mascaras_Lensless.pptx"

# Colores INTISAT (paleta sobria)
COLOR_TITLE = RGBColor(0x1A, 0x36, 0x5D)
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)
COLOR_ACCENT = RGBColor(0xC5, 0x8F, 0x00)
COLOR_OK = RGBColor(0x2F, 0x85, 0x5A)
COLOR_WARN = RGBColor(0xC5, 0x3C, 0x3C)
COLOR_GRID_HEADER = RGBColor(0x1A, 0x36, 0x5D)
COLOR_GRID_ROW1 = RGBColor(0xF7, 0xFA, 0xFC)
COLOR_GRID_ROW2 = RGBColor(0xFF, 0xFF, 0xFF)


# ── Helpers ───────────────────────────────────────────────────────────

def add_blank_slide(prs):
    """Agrega slide en blanco usando el layout BLANK del template."""
    blank = prs.slide_layouts[0]
    return prs.slides.add_slide(blank)


def add_title_box(slide, text, top=0.4, left=0.6, width=12.1, height=0.9,
                   font_size=26, bold=True, color=COLOR_TITLE):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    # Linea bajo el titulo
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top + height - 0.05),
                                    Inches(width), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()
    return box


def add_text_box(slide, text, top, left, width, height,
                  font_size=14, bold=False, color=COLOR_BODY,
                  align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return box


def add_bullet_list(slide, items, top, left, width, height,
                     font_size=14, bullet="•", color=COLOR_BODY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def add_image(slide, path, top, left, width=None, height=None):
    if width is not None:
        return slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                          width=Inches(width))
    if height is not None:
        return slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                          height=Inches(height))
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top))


def add_footer(slide, slide_num, total):
    add_text_box(slide,
                  "INTISAT Payload — Microscopia FPM lensless | Experimento de mascaras opticas",
                  top=7.1, left=0.6, width=10.5, height=0.3,
                  font_size=9, color=RGBColor(0x77, 0x77, 0x77))
    add_text_box(slide, f"{slide_num} / {total}",
                  top=7.1, left=12.1, width=0.7, height=0.3,
                  font_size=9, color=RGBColor(0x77, 0x77, 0x77),
                  align=PP_ALIGN.RIGHT)


def add_table(slide, data, top, left, width, height,
               header_color=COLOR_GRID_HEADER,
               row1_color=COLOR_GRID_ROW1, row2_color=COLOR_GRID_ROW2,
               font_size=11):
    """data: lista de listas (filas, primera fila = header)."""
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols,
                                          Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = tbl_shape.table
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.06)
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"
            if i == 0:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                run.font.color.rgb = COLOR_BODY
                cell.fill.solid()
                cell.fill.fore_color.rgb = row1_color if i % 2 else row2_color
    return tbl_shape


def add_colored_box(slide, text, top, left, width, height,
                      bg_color, border_color, title=None,
                      title_color=None, body_size=12):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    rect.line.color.rgb = border_color
    rect.line.width = Pt(1.2)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    if title:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = title_color or border_color
        run.font.name = "Calibri"
        if isinstance(text, list):
            for line in text:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.size = Pt(body_size)
                run.font.color.rgb = COLOR_BODY
                run.font.name = "Calibri"
        else:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = text
            run.font.size = Pt(body_size)
            run.font.color.rgb = COLOR_BODY
            run.font.name = "Calibri"
    return rect


# ── Construccion de slides ────────────────────────────────────────────

def build():
    prs = Presentation(str(TEMPLATE))
    # Eliminar slides demo del template (de atras hacia adelante)
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst)[::-1]:
        rId = sldId.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    SLIDES_TOTAL_PLACEHOLDER = 14

    # ─── Slide 1: PORTADA ────────────────────────────────────────────
    s1 = add_blank_slide(prs)
    # Banda superior decorativa
    band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0),
                                  Inches(13.33), Inches(0.5))
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_TITLE
    band.line.fill.background()
    add_text_box(s1, "INTISAT — Payload de microscopia",
                  top=0.05, left=0.6, width=12, height=0.4,
                  font_size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # Titulo principal
    add_text_box(s1, "Mascaras opticas para microscopia lensless",
                  top=2.0, left=0.8, width=11.7, height=1.1,
                  font_size=40, bold=True, color=COLOR_TITLE,
                  align=PP_ALIGN.CENTER)
    add_text_box(s1, "Estudio comparativo de rendijas: circular, cuadrada y lineal",
                  top=3.2, left=0.8, width=11.7, height=0.6,
                  font_size=22, color=RGBColor(0x55, 0x55, 0x55),
                  align=PP_ALIGN.CENTER)

    # Linea decorativa
    line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(5.0), Inches(4.1),
                                  Inches(3.33), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    add_text_box(s1, "Disenadas en CAD · Fabricacion + experimento + analisis",
                  top=4.4, left=0.8, width=11.7, height=0.4,
                  font_size=14, color=RGBColor(0x77, 0x77, 0x77),
                  align=PP_ALIGN.CENTER)

    # Datos al pie
    add_text_box(s1, "Sensor: OV5647 (lensless)  ·  Iluminacion: OLED SSD1351 128x128",
                  top=5.6, left=0.8, width=11.7, height=0.4,
                  font_size=14, bold=True, color=COLOR_TITLE,
                  align=PP_ALIGN.CENTER)
    add_text_box(s1, "CubeSat-EdgeAI-Payload  ·  github.com/JaminYC/CubeSat-EdgeAI-Payload",
                  top=6.0, left=0.8, width=11.7, height=0.4,
                  font_size=12, color=RGBColor(0x77, 0x77, 0x77),
                  align=PP_ALIGN.CENTER, font_name="Consolas")

    # Slide 2: INDICE
    s = add_blank_slide(prs)
    add_title_box(s, "Indice")
    items = [
        "1. Objetivo del experimento",
        "2. Concepto: por que una mascara fisica",
        "3. Las tres mascaras fabricadas (CAD)",
        "4. Tres modos de ensamblaje sobre el OV5647",
        "5. Hipotesis: rendija lineal (slit)",
        "6. Hipotesis: apertura cuadrada",
        "7. Hipotesis: pinhole circular",
        "8. Matriz experimental: 3 mascaras x 3 modos",
        "9. Metricas a medir y como compararlas",
        "10. Cuadro de resultados esperados",
        "11. Procedimiento paso a paso",
        "12. Riesgos y proximos pasos",
    ]
    add_bullet_list(s, items, top=1.5, left=1.5, width=10.5, height=5.0,
                     font_size=18)
    add_footer(s, 2, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 3: OBJETIVO
    s = add_blank_slide(prs)
    add_title_box(s, "Objetivo del experimento")
    add_text_box(s,
        "Caracterizar el efecto de tres geometrias de apertura sobre la calidad de imagen "
        "lensless del payload, montadas entre la muestra y el sensor OV5647.",
        top=1.6, left=0.8, width=11.7, height=1.0,
        font_size=16, color=COLOR_BODY)

    # Caja con tres preguntas que el experimento responde
    add_colored_box(s,
        ["¿Que mascara da el mejor compromiso entre nitidez y luz disponible?",
         "¿Que diferencias visibles aparecen entre rendija/cuadrada/circular?",
         "¿En que modo de ensamblaje (A/B/C) cada mascara funciona mejor?"],
        top=3.0, left=0.8, width=11.7, height=2.0,
        bg_color=RGBColor(0xFE, 0xF5, 0xE7), border_color=COLOR_ACCENT,
        title="Preguntas que el experimento responde:",
        title_color=COLOR_ACCENT, body_size=14)

    add_colored_box(s,
        ["3 mascaras · 3 modos = 9 condiciones experimentales",
         "Misma muestra (epidermis de cebolla) en todas las condiciones",
         "Misma iluminacion OLED (patron bright-field uniforme)",
         "Comparacion cuantitativa: contraste, resolucion, FOV utilizable"],
        top=5.2, left=0.8, width=11.7, height=1.4,
        bg_color=RGBColor(0xE6, 0xFF, 0xFA), border_color=COLOR_OK,
        title="Variables controladas:", title_color=COLOR_OK, body_size=13)

    add_footer(s, 3, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 4: CONCEPTO
    s = add_blank_slide(prs)
    add_title_box(s, "Concepto: por que necesitamos una mascara")
    add_text_box(s,
        "Sin lente, cada punto de la muestra dispersa luz hacia muchos pixeles del sensor → blur. "
        "La mascara restringe los angulos que llegan al sensor o a la muestra.",
        top=1.5, left=0.8, width=11.7, height=0.8,
        font_size=14, color=COLOR_BODY)
    # Formula
    add_colored_box(s,
        ["NA = sin(arctan(D / 2h))     ≈ D / (2h)  para angulos chicos",
         "",
         "donde:  D = tamaño caracteristico de la apertura  ·  h = distancia mascara - sensor"],
        top=2.5, left=0.8, width=11.7, height=1.4,
        bg_color=RGBColor(0xF7, 0xFA, 0xFC),
        border_color=COLOR_TITLE,
        title="Geometria fundamental:",
        body_size=14)
    # Tabla NA/h/D
    headers = ["h (mm)", "D = 0.5 mm", "D = 1.0 mm", "D = 2.0 mm"]
    rows = [
        ["0.5", "NA = 0.45 (26°)", "NA = 0.71 (45°)", "abierto"],
        ["1.0", "NA = 0.24 (14°)", "NA = 0.45 (26°)", "NA = 0.71 (45°)"],
        ["1.5", "NA = 0.16 (10°)", "NA = 0.32 (18°)", "NA = 0.55 (33°)"],
        ["2.0", "NA = 0.12 (7°)", "NA = 0.24 (14°)", "NA = 0.45 (26°)"],
    ]
    add_table(s, [headers] + rows,
               top=4.2, left=2.0, width=9.3, height=2.2, font_size=13)

    add_text_box(s,
        "Regla: NA bajo = nitidez alta + poca luz | NA alto = mucha luz + blur",
        top=6.55, left=0.8, width=11.7, height=0.4,
        font_size=13, color=COLOR_ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, 4, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 5: LAS TRES MASCARAS (CAD)
    s = add_blank_slide(prs)
    add_title_box(s, "Las tres mascaras fabricadas")
    add_text_box(s,
        "Cada mascara: placa exterior 25 x 25 mm con orejas de fijacion. "
        "Diferencia: SOLO la geometria del hueco central.",
        top=1.5, left=0.8, width=11.7, height=0.5,
        font_size=13, color=COLOR_BODY)

    # 3 imagenes lado a lado de las mascaras
    img_w = 4.0
    img_top = 2.2
    # Si los previews existen los ponemos, sino dejamos placeholders
    img_paths = [
        ROOT / "out" / "previews" / "rendija_(slit)_100_µm_x_4_mm.png",
        ROOT / "out" / "previews" / "apertura_cuadrada_1.5x1.5_mm.png",
        ROOT / "out" / "previews" / "pinhole_circular_d1.0_mm.png",
    ]
    labels = [
        "Rendija lineal (slit)",
        "Apertura cuadrada",
        "Pinhole circular"
    ]
    descs = [
        "100 μm × 4 mm",
        "1.5 × 1.5 mm",
        "Ø 1.0 mm"
    ]
    for i, (img, name, desc) in enumerate(zip(img_paths, labels, descs)):
        x = 0.5 + i * 4.3
        if img.exists():
            add_image(s, img, top=img_top, left=x, width=img_w)
        # Etiqueta
        add_text_box(s, name, top=img_top + 4.2, left=x, width=img_w, height=0.4,
                      font_size=14, bold=True, color=COLOR_TITLE,
                      align=PP_ALIGN.CENTER)
        add_text_box(s, desc, top=img_top + 4.55, left=x, width=img_w, height=0.4,
                      font_size=12, color=COLOR_ACCENT,
                      align=PP_ALIGN.CENTER, font_name="Consolas")

    add_footer(s, 5, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 6: MODOS DE ENSAMBLAJE
    s = add_blank_slide(prs)
    add_title_box(s, "Tres modos de ensamblaje sobre el OV5647")
    add_text_box(s,
        "La misma mascara puede montarse de tres formas diferentes. "
        "Cada modo tiene un compromiso distinto entre nitidez y control angular.",
        top=1.5, left=0.8, width=11.7, height=0.6,
        font_size=13, color=COLOR_BODY)
    asm = ROOT / "out" / "aperture_masks_assembly.png"
    if asm.exists():
        add_image(s, asm, top=2.2, left=0.4, width=12.5)
    add_footer(s, 6, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 7: HIPOTESIS - SLIT
    s = add_blank_slide(prs)
    add_title_box(s, "Hipotesis 1 — Rendija lineal (slit)")
    img = ROOT / "out" / "previews" / "rendija_(slit)_100_µm_x_4_mm.png"
    if img.exists():
        add_image(s, img, top=1.6, left=0.5, width=4.5)
    # Hipotesis a la derecha
    add_text_box(s, "Que esperamos ver",
                  top=1.7, left=5.5, width=7.5, height=0.5,
                  font_size=18, bold=True, color=COLOR_TITLE)
    add_bullet_list(s, [
        "Resolucion ALTA en el eje perpendicular a la rendija (eje corto: 100 μm)",
        "Resolucion BAJA en el eje paralelo a la rendija (campo abierto)",
        "Imagen tipo \"line scan\": una franja nitida y el resto borroso",
        "Util para confirmar el plano de enfoque y para line-scanning rotando la muestra",
        "Caso especial: con OLED puntual, da el mejor SNR de las tres",
    ], top=2.3, left=5.5, width=7.5, height=3.5, font_size=13)
    add_colored_box(s,
        ["NA_eje_corto = 100 μm / (2 × 1500 μm) ≈ 0.033",
         "Resolucion teorica (Abbe, λ=550 nm): d ≈ 8 μm en el eje fino",
         "Lo otro: campo abierto, blur libre"],
        top=5.5, left=0.5, width=12.3, height=1.4,
        bg_color=RGBColor(0xF7, 0xFA, 0xFC), border_color=COLOR_TITLE,
        title="Calculo:", body_size=12)
    add_footer(s, 7, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 8: HIPOTESIS - SQUARE
    s = add_blank_slide(prs)
    add_title_box(s, "Hipotesis 2 — Apertura cuadrada")
    img = ROOT / "out" / "previews" / "apertura_cuadrada_1.5x1.5_mm.png"
    if img.exists():
        add_image(s, img, top=1.6, left=0.5, width=4.5)
    add_text_box(s, "Que esperamos ver",
                  top=1.7, left=5.5, width=7.5, height=0.5,
                  font_size=18, bold=True, color=COLOR_TITLE)
    add_bullet_list(s, [
        "FOV cuadrado bien definido — bordes rectos visibles en el sensor",
        "Resolucion uniforme en X e Y (no privilegia ningun eje)",
        "Mas luz que el pinhole circular (area = 2.25 mm² vs 0.79 mm²)",
        "Diffraccion en las esquinas: aparecen patrones cruzados (sinc-cuadrado)",
        "Bueno para test de calibracion: el cuadrado en la imagen confirma alineacion",
    ], top=2.3, left=5.5, width=7.5, height=3.5, font_size=13)
    add_colored_box(s,
        ["NA = 0.75 mm / 1500 μm ≈ 0.50 en eje X o Y",
         "Diagonal: NA ≈ 0.71  →  resolucion direccional",
         "Transmitancia ≈ 32% (vs 11% del pinhole de 1 mm)"],
        top=5.5, left=0.5, width=12.3, height=1.4,
        bg_color=RGBColor(0xF7, 0xFA, 0xFC), border_color=COLOR_TITLE,
        title="Calculo:", body_size=12)
    add_footer(s, 8, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 9: HIPOTESIS - PINHOLE
    s = add_blank_slide(prs)
    add_title_box(s, "Hipotesis 3 — Pinhole circular")
    img = ROOT / "out" / "previews" / "pinhole_circular_d1.0_mm.png"
    if img.exists():
        add_image(s, img, top=1.6, left=0.5, width=4.5)
    add_text_box(s, "Que esperamos ver",
                  top=1.7, left=5.5, width=7.5, height=0.5,
                  font_size=18, bold=True, color=COLOR_TITLE)
    add_bullet_list(s, [
        "FOV circular — el contorno es un disco uniforme",
        "Resolucion isotropica (igual en todas direcciones)",
        "Patron Airy: anillos concentricos cuando hay difraccion",
        "El estandar de comparacion — todos los textos de optica usan pinhole circular",
        "Si quereis CONTAR celulas correctamente, este es el mas confiable",
    ], top=2.3, left=5.5, width=7.5, height=3.5, font_size=13)
    add_colored_box(s,
        ["NA = 0.5 mm / 1500 μm ≈ 0.33 (uniforme en todas direcciones)",
         "Resolucion Airy: 1.22·λ/(2·NA) ≈ 1.0 μm para λ=550 nm",
         "Transmitancia ≈ 11% (la mas baja de las tres)"],
        top=5.5, left=0.5, width=12.3, height=1.4,
        bg_color=RGBColor(0xF7, 0xFA, 0xFC), border_color=COLOR_TITLE,
        title="Calculo:", body_size=12)
    add_footer(s, 9, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 10: MATRIZ EXPERIMENTAL
    s = add_blank_slide(prs)
    add_title_box(s, "Matriz experimental: 3 mascaras × 3 modos = 9 condiciones")
    add_text_box(s,
        "Cada celda de la tabla es un experimento. Misma muestra, misma iluminacion.",
        top=1.5, left=0.8, width=11.7, height=0.4,
        font_size=13, color=COLOR_BODY)
    headers = ["Mascara \\ Modo", "A) contact", "B) lensless + spacer", "C) contact + mascara arriba"]
    rows = [
        ["Slit (rendija)",
         "shadow line scan + blur OLED",
         "line scan claro, perdida de luz",
         "DPC con corte X o Y"],
        ["Cuadrada",
         "FOV cuadrado borroso",
         "FOV cuadrado nitido, mas luz",
         "iluminacion cuadrada controlada"],
        ["Pinhole circular",
         "FOV circular borroso",
         "FOV circular nitido — referencia",
         "Köhler simplificado, contraste limpio"],
    ]
    add_table(s, [headers] + rows,
               top=2.2, left=0.5, width=12.3, height=2.5, font_size=12)

    add_colored_box(s,
        ["Total: 9 capturas (idealmente 3 repeticiones por celda → 27 capturas)",
         "Cada captura: 1 imagen RAW + 1 imagen procesada + log JSON con metadatos",
         "Tiempo estimado por celda: 5 min (incluye montaje y captura)",
         "Total experimento: ~ 1 - 2 h"],
        top=5.0, left=0.5, width=12.3, height=1.8,
        bg_color=RGBColor(0xE6, 0xFF, 0xFA), border_color=COLOR_OK,
        title="Plan logistico:", title_color=COLOR_OK, body_size=12)
    add_footer(s, 10, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 11: METRICAS
    s = add_blank_slide(prs)
    add_title_box(s, "Metricas a medir y como compararlas")
    add_text_box(s,
        "Para cada captura calculamos las siguientes metricas, ya implementadas en el pipeline.",
        top=1.5, left=0.8, width=11.7, height=0.5, font_size=13, color=COLOR_BODY)
    headers = ["Metrica", "Definicion", "Como se mide", "Que espero"]
    rows = [
        ["Resolucion (FWHM)",
         "Ancho a media altura de un borde",
         "Edge spread function en una pared celular",
         "Pinhole < cuadrada < slit (eje fino)"],
        ["Contraste (Michelson)",
         "(Imax - Imin) / (Imax + Imin)",
         "ROI sobre celula vs fondo",
         "Pinhole > slit > cuadrada"],
        ["SNR",
         "media / desvio en zona uniforme",
         "ROI lejos de la muestra",
         "Slit > cuadrada > pinhole"],
        ["FOV utilizable",
         "% del sensor con buena nitidez",
         "binarizar y contar pixels",
         "Cuadrada > pinhole > slit"],
        ["# celulas detectadas",
         "Salida del segmentador StarDist",
         "automatico via pipeline",
         "Cuadrada (mas FOV) o pinhole (mas nitido)"],
        ["Brillo medio",
         "valor de gris promedio",
         "media global ROI muestra",
         "Cuadrada > slit > pinhole"],
    ]
    add_table(s, [headers] + rows,
               top=2.2, left=0.4, width=12.5, height=4.5, font_size=11)
    add_footer(s, 11, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 12: RESULTADOS ESPERADOS
    s = add_blank_slide(prs)
    add_title_box(s, "Cuadro de resultados esperados")
    add_text_box(s,
        "Sintetiza nuestras hipotesis: que mascara tiene la ventaja en cada metrica.",
        top=1.5, left=0.8, width=11.7, height=0.4, font_size=13, color=COLOR_BODY)
    headers = ["", "Slit (rendija)", "Cuadrada", "Pinhole circular"]
    rows = [
        ["Resolucion (eje fino)", "★★★", "★★", "★★"],
        ["Resolucion uniforme", "—", "★★", "★★★"],
        ["Luz / SNR", "★★★", "★★", "★"],
        ["FOV", "★★ (linea)", "★★★", "★★"],
        ["Simplicidad de analisis", "★", "★★", "★★★"],
        ["Caso de uso", "line-scan, perfilometria", "lensless general", "calibracion, referencia"],
    ]
    add_table(s, [headers] + rows,
               top=2.1, left=0.5, width=12.3, height=3.8, font_size=13)

    add_colored_box(s,
        ["• Pinhole circular = referencia metrologica (la prediccion mas confiable).",
         "• Cuadrada = workhorse para uso general (mas luz, FOV completo).",
         "• Slit = caso especial; util para line-scanning o medir un eje con maxima resolucion."],
        top=6.1, left=0.5, width=12.3, height=1.0,
        bg_color=RGBColor(0xFE, 0xF5, 0xE7), border_color=COLOR_ACCENT,
        title="Recomendacion preliminar:", title_color=COLOR_ACCENT, body_size=12)
    add_footer(s, 12, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 13: PROCEDIMIENTO
    s = add_blank_slide(prs)
    add_title_box(s, "Procedimiento experimental paso a paso")
    pasos = [
        "1. Calibracion previa: regleta de 1 mm en cada modo → confirmar μm/px.",
        "2. Imprimir/cortar las 3 mascaras (PLA negro 0.4 mm o aluminio anodizado).",
        "3. Limpiar cover glass del OV5647 con alcohol isopropilico.",
        "4. Preparar muestra: corte de epidermis de cebolla (~ 5 mm × 5 mm).",
        "5. Para cada mascara × modo (9 combinaciones):",
        "    a. Montar la mascara con el spacer correspondiente (h = 1.5 o 4.0 mm).",
        "    b. Encender OLED en patron bright-field uniforme.",
        "    c. Capturar 3 imagenes (control de repetibilidad).",
        "    d. Guardar en /var/cubesat/incoming/exp_<mask>_<mode>_<rep>/.",
        "6. Procesar todas las capturas con pipeline.controller.run() → genera CSV de metricas.",
        "7. Comparar metricas (Slide 11) entre celdas → confirmar/refutar hipotesis.",
        "8. Documentar resultados en exp_<fecha>/RESULTS.md y subir a GitHub.",
    ]
    add_bullet_list(s, pasos, top=1.5, left=0.8, width=11.7, height=5.4,
                     font_size=13, bullet="")
    add_footer(s, 13, SLIDES_TOTAL_PLACEHOLDER)

    # Slide 14: RIESGOS Y NEXT
    s = add_blank_slide(prs)
    add_title_box(s, "Riesgos, limitaciones y proximos pasos")

    add_colored_box(s,
        ["• Tolerancias de fabricacion: si la rendija de 100 μm sale a 200 μm, la NA cambia ×2.",
         "• Alineacion mecanica: descentrado de 0.5 mm = sombra del borde de mascara.",
         "• Iluminacion no uniforme: el OLED puede tener gradiente → afecta contraste.",
         "• Reflexiones internas: el lado expuesto debe ser MATE (no brillante).",
         "• Limpieza: una mota en la mascara aparece como punto oscuro fijo en todas las imagenes."],
        top=1.5, left=0.5, width=12.3, height=2.4,
        bg_color=RGBColor(0xFE, 0xE2, 0xE2), border_color=COLOR_WARN,
        title="Riesgos identificados:", title_color=COLOR_WARN, body_size=12)

    add_colored_box(s,
        ["1. Una vez completado este experimento, repetir con OLED en single-LED (FPM) "
         "para ver si el blur cambia con iluminacion puntual.",
         "2. Probar mascara honeycomb (5×5 pinholes) como cuarto candidato.",
         "3. Implementar DPC con la mascara C en cuatro orientaciones.",
         "4. Medir profundidad de campo: muestra a distintas alturas.",
         "5. Integrar la mejor mascara al payload final del CubeSat."],
        top=4.1, left=0.5, width=12.3, height=2.6,
        bg_color=RGBColor(0xE6, 0xFF, 0xFA), border_color=COLOR_OK,
        title="Proximos pasos:", title_color=COLOR_OK, body_size=12)
    add_footer(s, 14, SLIDES_TOTAL_PLACEHOLDER)

    # Guardar
    prs.save(str(OUT_PATH))
    print(f"\n✓ Presentacion generada: {OUT_PATH}")
    print(f"  Slides totales: {len(prs.slides)}")


if __name__ == "__main__":
    build()
